# modules/pptx_export.py
from pptx import Presentation
from pptx.util import Inches
import os

def create_pptx(title, insights, image_paths, out_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = title
    for ins in insights:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Insight"
        tx = slide.shapes.placeholders[1].text_frame
        tx.text = ins
    for img in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(img, Inches(1), Inches(1.5), width=Inches(8))
    prs.save(out_path)
    return out_path